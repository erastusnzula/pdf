import os
import tempfile
import uuid

from django.http import HttpResponse
from django.shortcuts import render
from django.views import View
from django.conf import settings
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches

POPPLER_PATH = settings.POPPLER_PATH

TMP_DIR = os.path.join(os.getenv("TMP", "/tmp"))


def pdf_to_pptx(pdf_path, output_path):
    slides = Presentation()

    # Convert PDF → image pages
    pages = convert_from_path(pdf_path, dpi=200, poppler_path=POPPLER_PATH)

    for page in pages:
        slide = slides.slides.add_slide(slides.slide_layouts[6])  # blank slide

        # Save page temp
        page_path = f"{output_path}_{uuid.uuid4().hex}.png"
        page.save(page_path, "PNG")

        # Add image covering entire slide
        slide.shapes.add_picture(page_path,
                                 left=0,
                                 top=0,
                                 width=slides.slide_width,
                                 height=slides.slide_height)

        os.remove(page_path)  # delete temp image

    slides.save(output_path)


class PdfToPowerpoint(View):
    def get(self, *args, **kwargs):
        return render(self.request, 'pdf/convert_to_powerpoint.html')

    def post(self, *args, **kwargs):
        f = self.request.FILES.get("pdf")
        if not f:
            return HttpResponse("No PDF uploaded", status=400)

        name, ext = os.path.splitext(f.name)
        input_path = os.path.join(TMP_DIR, f"{uuid.uuid4().hex}{ext}")

        # Save PDF temporarily
        with open(input_path, "wb+") as fp:
            for chunk in f.chunks():
                fp.write(chunk)

        output_path = os.path.join(TMP_DIR, f"{name}.pptx")

        try:
            pdf_to_pptx(input_path, output_path)

            with open(output_path, "rb") as fp:
                response = HttpResponse(fp.read(),
                                        content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation")
                response["Content-Disposition"] = f'attachment; filename="{name}.pptx"'
                return response

        finally:
            try:
                os.remove(input_path)
            except:
                pass
            try:
                os.remove(output_path)
            except:
                pass
